# from __future__ import print_function
from pptx import Presentation

from datetime import date

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import git
from pathlib import Path

repo_root = git.Repo(".", search_parent_directories=True).working_dir


# Functions go here
class MakePowerpoint:
    def __init__(self):
        """
        Initializes a new instance of the PowerPoint class.
        """
        pass

    def create_title_slide(self, the_title):
        """
        Creates a title slide with the given title and today's date as the subtitle.

        Args:
            the_title (str): The title to be displayed on the slide.

        Returns:
            None
        """
        # Create a new PowerPoint presentation
        self.prs = Presentation()

        # Select the layout for the title slide
        title_slide_layout = self.prs.slide_layouts[0]

        # Add a new slide with the selected layout
        slide = self.prs.slides.add_slide(title_slide_layout)

        # Add a title to the slide
        title = slide.shapes.title
        title.text = the_title

        # Add today's date as the subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = "Generated on {:%m-%d-%Y}".format(date.today())

    def insert_slide(self, title_below, description, figure_save_loc):
        """
        Inserts a slide into the PowerPoint presentation with a title, description, and figure.

        Args:
            title_below (str): The title of the slide.
            description (str): The description of the slide.
            figure_save_loc (str): The file path of the figure to be added to the slide.

        Returns:
            None
        """
        # Select the slide layout to be used
        slide_layout = self.prs.slide_layouts[5]

        # Add a slide to the presentation using the selected layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Add a title to the slide
        title = slide.shapes.title
        title.text = title_below

        # Add a picture to the slide
        slide.shapes.add_picture(
            image_file=figure_save_loc,
            left=Inches(0.5),
            top=Inches(2),
            width=Inches(6),
            height=Inches(5),
        )

        # Add a text box to the slide
        textbox = slide.shapes.add_textbox(
            left=Inches(7), top=Inches(2), width=Inches(3), height=Inches(5)
        )

        # Add a text frame to the text box
        tf = textbox.text_frame

        # Add a paragraph to the text frame
        para = tf.add_paragraph()
        para.text = description
        para.alignment = PP_ALIGN.LEFT
        para.font.size = Pt(10)

    def save_ppt(self, pres_name):
        """
        Saves the PowerPoint presentation with the given name.

        Args:
            self: the PowerPoint presentation object
            pres_name (str): the name of the PowerPoint presentation to save

        Returns:
            None
        """
        self.prs.save(pres_name)

    def insert_slide_with_feed(
        self, title_below, description, figure_save_loc, feed_figure_save_loc
    ):
        """
        Inserts a slide with a title, two images, and a textbox with a description.

        Args:
            title_below (str): The title of the slide.
            description (str): The description to be added to the textbox.
            figure_save_loc (str): The file path of the first image to be added.
            feed_figure_save_loc (str): The file path of the second image to be added.

        Returns:
            None
        """
        # Select the slide layout to be used
        slide_layout = self.prs.slide_layouts[5]
        # Add a slide with the selected layout
        slide = self.prs.slides.add_slide(slide_layout)
        # Add a title to the slide
        title = slide.shapes.title
        title.text = title_below
        # Add the first image to the slide
        slide.shapes.add_picture(
            image_file=figure_save_loc,
            left=Inches(0.5),
            top=Inches(1),
            width=Inches(6),
            height=Inches(3),
        )
        # Add the second image to the slide
        slide.shapes.add_picture(
            image_file=feed_figure_save_loc,
            left=Inches(0.5),
            top=Inches(4),
            width=Inches(6),
            height=Inches(3),
        )
        # Add a textbox to the slide
        textbox = slide.shapes.add_textbox(
            left=Inches(7), top=Inches(2), width=Inches(3), height=Inches(5)
        )
        # Add a paragraph to the textbox
        tf = textbox.text_frame
        para = tf.add_paragraph()
        para.text = description
        para.alignment = PP_ALIGN.LEFT
        para.font.size = Pt(10)

    def save_ppt(self, pres_name):
        self.prs.save(pres_name)


if __name__ == "__main__":
    the_title = "the_title"
    mp = MakePowerpoint()
    mp.create_title_slide(the_title)

    mp.insert_slide(
        "title_below",
        "description",
        Path(repo_root) / "results" / "large_reports" / the_title / ".png",
    )
    mp.save_ppt("baseline.pptx")
